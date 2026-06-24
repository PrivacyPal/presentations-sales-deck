#!/usr/bin/env python3
"""Assemble slide JPEGs into a native PowerPoint (.pptx), one full-bleed
image per 16:9 widescreen slide. Each slide IS the rendered image, so there
is no font dependency, no reflow, and no pagination drift -- identical to the
PDF output, but as an editable/presentable .pptx.

Usage:
    python3 build_pptx.py <jpg_dir> <out.pptx>
    python3 build_pptx.py                      # defaults to the SME deck
"""
import sys, glob, os
from pptx import Presentation
from pptx.util import Inches

BASE = "/Users/jason/dev/PrivacyPal/privacypal-sales-deck"
jpg_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.join(BASE, "slides_jpg")
out     = sys.argv[2] if len(sys.argv) > 2 else os.path.join(BASE, "privacypal-sales-deck.pptx")

# PowerPoint widescreen: 13.333in x 7.5in (matches the 1280x720 / 16:9 stage)
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
blank = prs.slide_layouts[6]  # fully blank layout

jpgs = sorted(glob.glob(os.path.join(jpg_dir, "slide-*.jpg")))
if not jpgs:
    sys.exit("no jpgs found in %s" % jpg_dir)

for path in jpgs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(path, 0, 0, width=EMU_W, height=EMU_H)

prs.save(out)
print("wrote %s : %d slides" % (out, len(jpgs)))
